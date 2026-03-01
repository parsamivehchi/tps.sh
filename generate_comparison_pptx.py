#!/usr/bin/env python3
"""
Generate a polished comparison PowerPoint presentation for tps.sh results.
Run ID: 20260226_002903 | 84 tests | 4 local Ollama models x 21 prompts x 7 categories
Hardware: Apple M2 Max, 32GB RAM
"""

import os
import io
import tempfile
from pathlib import Path

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Constants & palette
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

BG_DARK = "#0f172a"       # slate-900
BG_CARD = "#1e293b"       # slate-800
BG_CARD_ALT = "#1a2332"
TEXT_WHITE = "#f8fafc"
TEXT_MUTED = "#94a3b8"
ACCENT_LINE = "#334155"

MODEL_COLORS = {
    "qwen3-coder":      "#10b981",
    "qwen2.5-coder":    "#3b82f6",
    "deepseek-r1":      "#f59e0b",
    "glm-4.7-flash":    "#f43f5e",
}

MODEL_LABELS = {
    "qwen3-coder":      "Qwen3-Coder (30B MoE)",
    "qwen2.5-coder":    "Qwen2.5-Coder (14B)",
    "deepseek-r1":      "DeepSeek-R1 (14B)",
    "glm-4.7-flash":    "GLM-4.7-Flash (~9B)",
}

MODEL_ORDER = ["qwen3-coder", "qwen2.5-coder", "deepseek-r1", "glm-4.7-flash"]

CATEGORIES = [
    "code_generation",
    "debugging_reasoning",
    "explanation_teaching",
    "long_complex",
    "refactoring",
    "short_quick",
    "tool_calling",
]

CAT_LABELS = {
    "code_generation": "Code Gen",
    "debugging_reasoning": "Debug/Reason",
    "explanation_teaching": "Explain/Teach",
    "long_complex": "Long/Complex",
    "refactoring": "Refactoring",
    "short_quick": "Short/Quick",
    "tool_calling": "Tool Calling",
}

# Category TPS data
CAT_TPS = {
    "code_generation":      {"qwen3-coder": 48.4, "qwen2.5-coder": 15.7, "deepseek-r1": 15.0, "glm-4.7-flash": 10.4},
    "debugging_reasoning":  {"qwen3-coder": 49.9, "qwen2.5-coder": 15.8, "deepseek-r1": 14.6, "glm-4.7-flash": 8.0},
    "explanation_teaching": {"qwen3-coder": 43.7, "qwen2.5-coder": 15.6, "deepseek-r1": 14.6, "glm-4.7-flash": 0.0},
    "long_complex":         {"qwen3-coder": 45.1, "qwen2.5-coder": 15.7, "deepseek-r1": 14.8, "glm-4.7-flash": 4.5},
    "refactoring":          {"qwen3-coder": 49.0, "qwen2.5-coder": 15.5, "deepseek-r1": 13.7, "glm-4.7-flash": 14.0},
    "short_quick":          {"qwen3-coder": 54.5, "qwen2.5-coder": 15.9, "deepseek-r1": 14.3, "glm-4.7-flash": 12.4},
    "tool_calling":         {"qwen3-coder": 50.9, "qwen2.5-coder": 15.2, "deepseek-r1": 15.1, "glm-4.7-flash": 21.9},
}

# Aggregate model data
MODEL_DATA = {
    "qwen3-coder":   {"tps_avg": 48.8, "tps_min": 43, "tps_max": 59, "tps_std": None, "ttft": 1064, "total_time": 37.8, "tokens": 35063, "vram": 19449},
    "qwen2.5-coder": {"tps_avg": 15.6, "tps_min": 14.5, "tps_max": 16.3, "tps_std": None, "ttft": 1513, "total_time": 68.2, "tokens": 21595, "vram": 13864},
    "deepseek-r1":   {"tps_avg": 14.6, "tps_min": 13.2, "tps_max": 16.4, "tps_std": None, "ttft": 70171, "total_time": 137.0, "tokens": 40662, "vram": 13864},
    "glm-4.7-flash": {"tps_avg": 10.2, "tps_min": 0, "tps_max": 28.7, "tps_std": 8.8, "ttft": 54815, "total_time": 229.5, "tokens": 33176, "vram": 19991},
}

# Compute rough std for models that didn't provide one
for k, v in MODEL_DATA.items():
    if v["tps_std"] is None:
        v["tps_std"] = (v["tps_max"] - v["tps_min"]) / 4  # rough estimate

OUTPUT_DIR = Path("/Users/parsamivehchi/Desktop/DEV/LLM-BENCH/reports")
OUTPUT_FILE = OUTPUT_DIR / "model_comparison_20260226.pptx"
CHART_DIR = Path(tempfile.mkdtemp(prefix="llmbench_charts_"))


# ---------------------------------------------------------------------------
# Matplotlib helpers
# ---------------------------------------------------------------------------
def hex_to_rgb_tuple(h: str):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) / 255.0 for i in (0, 2, 4))


def setup_dark_style():
    """Return a dict of rcParams for dark-themed charts."""
    return {
        "figure.facecolor": BG_DARK,
        "axes.facecolor": BG_CARD,
        "axes.edgecolor": ACCENT_LINE,
        "axes.labelcolor": TEXT_WHITE,
        "text.color": TEXT_WHITE,
        "xtick.color": TEXT_MUTED,
        "ytick.color": TEXT_MUTED,
        "grid.color": ACCENT_LINE,
        "grid.alpha": 0.5,
        "legend.facecolor": BG_CARD,
        "legend.edgecolor": ACCENT_LINE,
        "legend.labelcolor": TEXT_WHITE,
        "font.family": "sans-serif",
        "font.size": 14,
    }


def save_chart(fig, name: str) -> str:
    path = str(CHART_DIR / f"{name}.png")
    fig.savefig(path, dpi=200, bbox_inches="tight", facecolor=fig.get_facecolor(), pad_inches=0.3)
    plt.close(fig)
    return path


# ---------------------------------------------------------------------------
# Chart generators
# ---------------------------------------------------------------------------
def chart_speed_comparison() -> str:
    """Horizontal bar chart of average TPS."""
    with plt.rc_context(setup_dark_style()):
        fig, ax = plt.subplots(figsize=(10, 4.5))
        models = MODEL_ORDER[::-1]
        tps = [MODEL_DATA[m]["tps_avg"] for m in models]
        colors = [MODEL_COLORS[m] for m in models]
        labels = [MODEL_LABELS[m] for m in models]

        bars = ax.barh(labels, tps, color=colors, height=0.55, edgecolor="none", zorder=3)
        ax.set_xlabel("Tokens per Second (avg)", fontsize=13, fontweight="bold")
        ax.set_xlim(0, max(tps) * 1.2)
        ax.grid(axis="x", linestyle="--", alpha=0.3, zorder=0)
        ax.tick_params(axis="y", labelsize=12)

        for bar, val in zip(bars, tps):
            ax.text(bar.get_width() + 0.8, bar.get_y() + bar.get_height() / 2,
                    f" {val:.1f} tok/s", va="center", ha="left", fontsize=13,
                    fontweight="bold", color=TEXT_WHITE)

        fig.suptitle("Generation Speed Comparison", fontsize=18, fontweight="bold", y=0.98, color=TEXT_WHITE)
        fig.tight_layout(rect=[0, 0, 1, 0.93])
    return save_chart(fig, "speed_comparison")


def chart_ttft_comparison() -> str:
    """Bar chart of Time To First Token."""
    with plt.rc_context(setup_dark_style()):
        fig, ax = plt.subplots(figsize=(10, 4.5))
        models = MODEL_ORDER
        ttft = [MODEL_DATA[m]["ttft"] / 1000.0 for m in models]  # convert to seconds
        colors = [MODEL_COLORS[m] for m in models]
        labels = [MODEL_LABELS[m] for m in models]

        bars = ax.bar(labels, ttft, color=colors, width=0.55, edgecolor="none", zorder=3)
        ax.set_ylabel("Time to First Token (seconds)", fontsize=13, fontweight="bold")
        ax.grid(axis="y", linestyle="--", alpha=0.3, zorder=0)
        ax.tick_params(axis="x", labelsize=10)

        for bar, val in zip(bars, ttft):
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.5,
                    f"{val:.1f}s", ha="center", va="bottom", fontsize=12,
                    fontweight="bold", color=TEXT_WHITE)

        ax.set_ylim(0, max(ttft) * 1.18)
        fig.suptitle("Latency: Time to First Token", fontsize=18, fontweight="bold", y=0.98, color=TEXT_WHITE)
        fig.tight_layout(rect=[0, 0, 1, 0.93])
    return save_chart(fig, "ttft_comparison")


def chart_response_time() -> str:
    """Bar chart of average total response time."""
    with plt.rc_context(setup_dark_style()):
        fig, ax = plt.subplots(figsize=(10, 4.5))
        models = MODEL_ORDER
        times = [MODEL_DATA[m]["total_time"] for m in models]
        colors = [MODEL_COLORS[m] for m in models]
        labels = [MODEL_LABELS[m] for m in models]

        bars = ax.bar(labels, times, color=colors, width=0.55, edgecolor="none", zorder=3)
        ax.set_ylabel("Average Response Time (seconds)", fontsize=13, fontweight="bold")
        ax.grid(axis="y", linestyle="--", alpha=0.3, zorder=0)
        ax.tick_params(axis="x", labelsize=10)

        for bar, val in zip(bars, times):
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1,
                    f"{val:.1f}s", ha="center", va="bottom", fontsize=12,
                    fontweight="bold", color=TEXT_WHITE)

        ax.set_ylim(0, max(times) * 1.15)
        fig.suptitle("Average Total Response Time", fontsize=18, fontweight="bold", y=0.98, color=TEXT_WHITE)
        fig.tight_layout(rect=[0, 0, 1, 0.93])
    return save_chart(fig, "response_time")


def chart_category_heatmap() -> str:
    """Heatmap grid of TPS by category and model."""
    with plt.rc_context(setup_dark_style()):
        fig, ax = plt.subplots(figsize=(11, 5.5))

        data = []
        for cat in CATEGORIES:
            row = [CAT_TPS[cat][m] for m in MODEL_ORDER]
            data.append(row)
        data = np.array(data)

        im = ax.imshow(data, cmap="YlOrRd", aspect="auto", vmin=0, vmax=60)

        ax.set_xticks(range(len(MODEL_ORDER)))
        ax.set_xticklabels([MODEL_LABELS[m] for m in MODEL_ORDER], fontsize=10, rotation=15, ha="right")
        ax.set_yticks(range(len(CATEGORIES)))
        ax.set_yticklabels([CAT_LABELS[c] for c in CATEGORIES], fontsize=11)

        # Annotate cells
        for i in range(len(CATEGORIES)):
            for j in range(len(MODEL_ORDER)):
                val = data[i, j]
                text_color = "#0f172a" if val > 25 else TEXT_WHITE
                ax.text(j, i, f"{val:.1f}", ha="center", va="center",
                        fontsize=13, fontweight="bold", color=text_color)

        cbar = fig.colorbar(im, ax=ax, shrink=0.8, pad=0.02)
        cbar.set_label("Tokens/sec", fontsize=12, color=TEXT_WHITE)
        cbar.ax.yaxis.set_tick_params(color=TEXT_MUTED)
        plt.setp(cbar.ax.yaxis.get_ticklabels(), color=TEXT_MUTED)

        fig.suptitle("Category Performance Heatmap (tok/s)", fontsize=18, fontweight="bold", y=0.98, color=TEXT_WHITE)
        fig.tight_layout(rect=[0, 0, 1, 0.93])
    return save_chart(fig, "category_heatmap")


def chart_tokens() -> str:
    """Bar chart of total tokens generated."""
    with plt.rc_context(setup_dark_style()):
        fig, ax = plt.subplots(figsize=(10, 4.5))
        models = MODEL_ORDER
        tokens = [MODEL_DATA[m]["tokens"] for m in models]
        colors = [MODEL_COLORS[m] for m in models]
        labels = [MODEL_LABELS[m] for m in models]

        bars = ax.bar(labels, tokens, color=colors, width=0.55, edgecolor="none", zorder=3)
        ax.set_ylabel("Total Tokens Generated", fontsize=13, fontweight="bold")
        ax.grid(axis="y", linestyle="--", alpha=0.3, zorder=0)
        ax.tick_params(axis="x", labelsize=10)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x/1000:.0f}K"))

        for bar, val in zip(bars, tokens):
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 300,
                    f"{val:,}", ha="center", va="bottom", fontsize=12,
                    fontweight="bold", color=TEXT_WHITE)

        ax.set_ylim(0, max(tokens) * 1.15)
        fig.suptitle("Output Verbosity: Total Tokens (84 tests)", fontsize=18, fontweight="bold", y=0.98, color=TEXT_WHITE)
        fig.tight_layout(rect=[0, 0, 1, 0.93])
    return save_chart(fig, "tokens")


def chart_vram() -> str:
    """Horizontal bar chart of VRAM usage."""
    with plt.rc_context(setup_dark_style()):
        fig, ax = plt.subplots(figsize=(10, 4.5))
        models = MODEL_ORDER[::-1]
        vram = [MODEL_DATA[m]["vram"] / 1024.0 for m in models]  # GB
        colors = [MODEL_COLORS[m] for m in models]
        labels = [MODEL_LABELS[m] for m in models]

        bars = ax.barh(labels, vram, color=colors, height=0.55, edgecolor="none", zorder=3)
        ax.set_xlabel("VRAM Usage (GB)", fontsize=13, fontweight="bold")
        ax.grid(axis="x", linestyle="--", alpha=0.3, zorder=0)
        ax.tick_params(axis="y", labelsize=12)

        for bar, val in zip(bars, vram):
            ax.text(bar.get_width() + 0.2, bar.get_y() + bar.get_height() / 2,
                    f" {val:.1f} GB", va="center", ha="left", fontsize=13,
                    fontweight="bold", color=TEXT_WHITE)

        ax.set_xlim(0, max(vram) * 1.2)
        fig.suptitle("VRAM Consumption", fontsize=18, fontweight="bold", y=0.98, color=TEXT_WHITE)
        fig.tight_layout(rect=[0, 0, 1, 0.93])
    return save_chart(fig, "vram")


def chart_consistency() -> str:
    """Range plot showing min/max/avg TPS with error bars."""
    with plt.rc_context(setup_dark_style()):
        fig, ax = plt.subplots(figsize=(10, 5))
        models = MODEL_ORDER
        labels = [MODEL_LABELS[m] for m in models]
        colors = [MODEL_COLORS[m] for m in models]

        y_pos = range(len(models))
        avgs = [MODEL_DATA[m]["tps_avg"] for m in models]
        mins = [MODEL_DATA[m]["tps_min"] for m in models]
        maxs = [MODEL_DATA[m]["tps_max"] for m in models]
        stds = [MODEL_DATA[m]["tps_std"] for m in models]

        for i, m in enumerate(models):
            color = MODEL_COLORS[m]
            # range line
            ax.plot([mins[i], maxs[i]], [i, i], color=color, linewidth=3, alpha=0.5, zorder=2)
            # min/max markers
            ax.scatter([mins[i], maxs[i]], [i, i], color=color, s=80, zorder=3, edgecolors="white", linewidths=0.5)
            # avg marker
            ax.scatter([avgs[i]], [i], color=color, s=200, zorder=4, marker="D", edgecolors="white", linewidths=1.5)
            # std annotation
            ax.annotate(f"avg={avgs[i]:.1f}  std={stds[i]:.1f}",
                        xy=(avgs[i], i), xytext=(avgs[i] + 3, i + 0.25),
                        fontsize=11, fontweight="bold", color=color)

        ax.set_yticks(y_pos)
        ax.set_yticklabels(labels, fontsize=12)
        ax.set_xlabel("Tokens per Second", fontsize=13, fontweight="bold")
        ax.grid(axis="x", linestyle="--", alpha=0.3, zorder=0)
        ax.invert_yaxis()

        # legend
        ax.scatter([], [], marker="D", s=100, color=TEXT_MUTED, label="Average TPS")
        ax.scatter([], [], marker="o", s=60, color=TEXT_MUTED, label="Min / Max")
        ax.plot([], [], linewidth=3, color=TEXT_MUTED, alpha=0.5, label="Range")
        ax.legend(loc="lower right", fontsize=10, framealpha=0.8)

        fig.suptitle("Consistency Analysis: TPS Range & Variance", fontsize=18, fontweight="bold", y=0.98, color=TEXT_WHITE)
        fig.tight_layout(rect=[0, 0, 1, 0.93])
    return save_chart(fig, "consistency")


def chart_head_to_head() -> str:
    """Radar chart comparing qwen3-coder to each model across categories."""
    with plt.rc_context(setup_dark_style()):
        fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(polar=True))
        fig.patch.set_facecolor(BG_DARK)
        ax.set_facecolor(BG_CARD)

        cats = CATEGORIES
        N = len(cats)
        angles = np.linspace(0, 2 * np.pi, N, endpoint=False).tolist()
        angles += angles[:1]

        for m in MODEL_ORDER:
            vals = [CAT_TPS[c][m] for c in cats]
            vals += vals[:1]
            color = MODEL_COLORS[m]
            lw = 3 if m == "qwen3-coder" else 1.8
            alpha_fill = 0.15 if m == "qwen3-coder" else 0.05
            ax.plot(angles, vals, color=color, linewidth=lw, label=MODEL_LABELS[m])
            ax.fill(angles, vals, color=color, alpha=alpha_fill)

        ax.set_xticks(angles[:-1])
        ax.set_xticklabels([CAT_LABELS[c] for c in cats], fontsize=10, color=TEXT_WHITE)
        ax.set_ylim(0, 60)
        ax.set_yticks([10, 20, 30, 40, 50])
        ax.set_yticklabels(["10", "20", "30", "40", "50"], fontsize=9, color=TEXT_MUTED)
        ax.yaxis.grid(color=ACCENT_LINE, alpha=0.4)
        ax.xaxis.grid(color=ACCENT_LINE, alpha=0.4)
        ax.spines["polar"].set_color(ACCENT_LINE)

        ax.legend(loc="upper right", bbox_to_anchor=(1.35, 1.12), fontsize=10, framealpha=0.9)
        fig.suptitle("Head-to-Head: Qwen3-Coder vs All", fontsize=18, fontweight="bold", y=0.97, color=TEXT_WHITE)
        fig.tight_layout(rect=[0, 0, 0.95, 0.93])
    return save_chart(fig, "head_to_head")


# ---------------------------------------------------------------------------
# PowerPoint helpers
# ---------------------------------------------------------------------------
def hex_to_rgbcolor(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_bg(slide, color_hex: str):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgbcolor(color_hex)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=TEXT_WHITE, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = hex_to_rgbcolor(color)
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=16, color=TEXT_WHITE, line_spacing=1.2,
                          bold=False, font_name="Calibri", alignment=PP_ALIGN.LEFT):
    """Add a textbox with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line_data, dict):
            p.text = line_data.get("text", "")
            p.font.size = Pt(line_data.get("size", font_size))
            p.font.bold = line_data.get("bold", bold)
            p.font.color.rgb = hex_to_rgbcolor(line_data.get("color", color))
            p.font.name = line_data.get("font", font_name)
            p.alignment = line_data.get("alignment", alignment)
        else:
            p.text = str(line_data)
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = hex_to_rgbcolor(color)
            p.font.name = font_name
            p.alignment = alignment

        p.space_after = Pt(font_size * (line_spacing - 1) + 2)

    return txBox


def add_shape_rect(slide, left, top, width, height, fill_hex):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgbcolor(fill_hex)
    shape.line.fill.background()
    return shape


def add_accent_bar(slide, left, top, width, color_hex):
    """Thin colored accent bar."""
    return add_shape_rect(slide, left, top, width, Pt(4), color_hex)


def add_chart_image(slide, chart_path, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(chart_path, left, top, width, height)
    else:
        slide.shapes.add_picture(chart_path, left, top, width=width)


def add_slide_number(slide, num, total):
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(1), Inches(0.4),
                f"{num}/{total}", font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def build_title_slide(prs, slide_num, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # accent bar at top
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), "#10b981")

    # title
    add_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(1),
                "tps.sh: Local Model Comparison", font_size=40, bold=True, color=TEXT_WHITE)

    # subtitle
    add_textbox(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.6),
                "Performance Benchmarking of 4 Ollama Models on Apple M2 Max",
                font_size=22, color=TEXT_MUTED)

    add_accent_bar(slide, Inches(1), Inches(3.4), Inches(3), "#10b981")

    # metadata
    meta_lines = [
        {"text": f"Run ID:  20260226_002903", "size": 16, "color": TEXT_MUTED, "bold": False},
        {"text": f"Tests:  84 (4 models x 21 prompts x 7 categories)", "size": 16, "color": TEXT_MUTED},
        {"text": f"Hardware:  Apple M2 Max  |  32 GB Unified RAM", "size": 16, "color": TEXT_MUTED},
        {"text": f"Runtime:  2h 48m", "size": 16, "color": TEXT_MUTED},
        {"text": f"Date:  February 26, 2026", "size": 16, "color": TEXT_MUTED},
    ]
    add_multiline_textbox(slide, Inches(1), Inches(3.8), Inches(6), Inches(3),
                          meta_lines, line_spacing=1.5)

    # Model color legend on the right
    add_textbox(slide, Inches(8), Inches(3.8), Inches(4), Inches(0.5),
                "Models Tested", font_size=18, bold=True, color=TEXT_WHITE)

    y_start = 4.5
    for m in MODEL_ORDER:
        color = MODEL_COLORS[m]
        add_shape_rect(slide, Inches(8), Inches(y_start), Inches(0.35), Inches(0.35), color)
        add_textbox(slide, Inches(8.55), Inches(y_start - 0.03), Inches(4), Inches(0.4),
                    MODEL_LABELS[m], font_size=15, color=TEXT_WHITE)
        y_start += 0.55

    add_slide_number(slide, slide_num, total_slides)
    return slide


def build_executive_summary(prs, slide_num, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.8),
                "Executive Summary", font_size=32, bold=True)
    add_accent_bar(slide, Inches(0.7), Inches(1.15), Inches(2.5), "#10b981")

    findings = [
        {
            "title": "Qwen3-Coder dominates speed",
            "detail": "At 48.8 tok/s average, it is 3.1x faster than the next model and maintains consistent performance across all 7 categories.",
            "color": MODEL_COLORS["qwen3-coder"],
        },
        {
            "title": "DeepSeek-R1 is the most verbose",
            "detail": "Generates 40,662 tokens total (most of any model) with 70s TTFT due to its chain-of-thought reasoning process.",
            "color": MODEL_COLORS["deepseek-r1"],
        },
        {
            "title": "GLM-4.7-Flash is the most inconsistent",
            "detail": "TPS ranges from 0 to 28.7 (std=8.8). Failed entirely on explanation/teaching tasks. Highest VRAM despite smallest parameter count.",
            "color": MODEL_COLORS["glm-4.7-flash"],
        },
        {
            "title": "Qwen2.5-Coder is the most efficient",
            "detail": "Lowest VRAM at 13.5 GB, steady 15.6 tok/s with minimal variance (14.5-16.3), and fewest total tokens (21,595).",
            "color": MODEL_COLORS["qwen2.5-coder"],
        },
    ]

    y = 1.6
    for f in findings:
        # colored left bar
        add_shape_rect(slide, Inches(0.7), Inches(y), Pt(5), Inches(0.9), f["color"])
        # card background
        add_shape_rect(slide, Inches(0.85), Inches(y), Inches(11.5), Inches(0.9), BG_CARD)
        add_textbox(slide, Inches(1.1), Inches(y + 0.05), Inches(11), Inches(0.4),
                    f["title"], font_size=18, bold=True, color=f["color"])
        add_textbox(slide, Inches(1.1), Inches(y + 0.45), Inches(11), Inches(0.45),
                    f["detail"], font_size=14, color=TEXT_MUTED)
        y += 1.15

    # Quick stats row at bottom
    stats = [
        ("48.8 tok/s", "Top Speed", "#10b981"),
        ("1.06s", "Best TTFT", "#10b981"),
        ("70.2s", "Worst TTFT", "#f59e0b"),
        ("40,662", "Max Tokens", "#f59e0b"),
        ("~20 GB", "Max VRAM", "#f43f5e"),
    ]
    x_start = 0.7
    for val, label, color in stats:
        add_shape_rect(slide, Inches(x_start), Inches(6.3), Inches(2.2), Inches(0.9), BG_CARD)
        add_textbox(slide, Inches(x_start), Inches(6.3), Inches(2.2), Inches(0.5),
                    val, font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x_start), Inches(6.75), Inches(2.2), Inches(0.35),
                    label, font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
        x_start += 2.45

    add_slide_number(slide, slide_num, total_slides)
    return slide


def build_chart_slide(prs, title, chart_path, slide_num, total_slides, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11), Inches(0.7),
                title, font_size=28, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(0.95), Inches(11), Inches(0.4),
                    subtitle, font_size=14, color=TEXT_MUTED)

    add_accent_bar(slide, Inches(0.7), Inches(1.0) if not subtitle else Inches(1.35),
                   Inches(2), "#10b981")

    chart_top = Inches(1.3) if not subtitle else Inches(1.6)
    add_chart_image(slide, chart_path, Inches(0.7), chart_top, Inches(11.8), Inches(5.6))

    add_slide_number(slide, slide_num, total_slides)
    return slide


def build_takeaways_slide(prs, slide_num, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.8),
                "Key Takeaways & Recommendations", font_size=32, bold=True)
    add_accent_bar(slide, Inches(0.7), Inches(1.15), Inches(2.5), "#10b981")

    takeaways = [
        {
            "heading": "Best Overall: Qwen3-Coder (30B MoE)",
            "points": [
                "3.1x faster than any other model tested (48.8 vs 15.6 tok/s)",
                "Consistent performance across all categories (43-59 tok/s range)",
                "Fastest time-to-first-token at 1.06 seconds",
                "MoE architecture delivers excellent speed despite 30B parameters",
            ],
            "color": "#10b981",
        },
        {
            "heading": "Best Efficiency: Qwen2.5-Coder (14B)",
            "points": [
                "Lowest VRAM footprint (13.5 GB) -- ideal for memory-constrained setups",
                "Most predictable output: minimal variance (14.5-16.3 tok/s)",
                "Most concise responses (21,595 total tokens across 84 tests)",
            ],
            "color": "#3b82f6",
        },
        {
            "heading": "Use With Caution: DeepSeek-R1 & GLM-4.7",
            "points": [
                "DeepSeek-R1: 70s TTFT makes it impractical for interactive use; best for async batch tasks",
                "GLM-4.7-Flash: Failed on explanation tasks; 0-28.7 tok/s range is unreliable; highest VRAM despite smallest model",
            ],
            "color": "#f59e0b",
        },
    ]

    y = 1.5
    for t in takeaways:
        add_shape_rect(slide, Inches(0.7), Inches(y), Pt(5), Inches(len(t["points"]) * 0.38 + 0.55), t["color"])
        add_shape_rect(slide, Inches(0.85), Inches(y), Inches(11.5), Inches(len(t["points"]) * 0.38 + 0.55), BG_CARD)

        add_textbox(slide, Inches(1.1), Inches(y + 0.08), Inches(11), Inches(0.4),
                    t["heading"], font_size=18, bold=True, color=t["color"])

        bullet_y = y + 0.5
        for pt in t["points"]:
            add_textbox(slide, Inches(1.3), Inches(bullet_y), Inches(10.8), Inches(0.35),
                        f"\u2022  {pt}", font_size=13, color=TEXT_MUTED)
            bullet_y += 0.38

        y += len(t["points"]) * 0.38 + 0.75

    # Footer
    add_textbox(slide, Inches(0.7), Inches(6.8), Inches(11), Inches(0.5),
                "tps.sh  |  Run 20260226_002903  |  Apple M2 Max 32GB  |  84 tests across 7 categories",
                font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, slide_num, total_slides)
    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    print("Generating charts...")
    chart_speed = chart_speed_comparison()
    chart_ttft = chart_ttft_comparison()
    chart_resp = chart_response_time()
    chart_cat = chart_category_heatmap()
    chart_tok = chart_tokens()
    chart_v = chart_vram()
    chart_con = chart_consistency()
    chart_h2h = chart_head_to_head()
    print(f"  Charts saved to {CHART_DIR}")

    print("Building presentation...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    total_slides = 11

    # 1. Title
    build_title_slide(prs, 1, total_slides)

    # 2. Executive Summary
    build_executive_summary(prs, 2, total_slides)

    # 3. Speed Comparison
    build_chart_slide(prs, "Generation Speed Comparison",
                      chart_speed, 3, total_slides,
                      subtitle="Average tokens per second across all 21 prompts per model")

    # 4. Latency (TTFT)
    build_chart_slide(prs, "Latency: Time to First Token",
                      chart_ttft, 4, total_slides,
                      subtitle="How long until the model starts generating? DeepSeek-R1 and GLM think extensively before responding.")

    # 5. Response Time
    build_chart_slide(prs, "Average Total Response Time",
                      chart_resp, 5, total_slides,
                      subtitle="End-to-end time from prompt submission to completion (seconds)")

    # 6. Category Heatmap
    build_chart_slide(prs, "Category Performance Breakdown",
                      chart_cat, 6, total_slides,
                      subtitle="Tokens/sec by task category -- GLM-4.7 failed explanation/teaching entirely (0 tok/s)")

    # 7. Output Verbosity
    build_chart_slide(prs, "Output Verbosity: Total Tokens Generated",
                      chart_tok, 7, total_slides,
                      subtitle="DeepSeek-R1's chain-of-thought produces 88% more tokens than Qwen2.5-Coder")

    # 8. VRAM Usage
    build_chart_slide(prs, "VRAM Consumption",
                      chart_v, 8, total_slides,
                      subtitle="Memory footprint during inference -- GLM uses the most VRAM despite being the smallest model")

    # 9. Consistency
    build_chart_slide(prs, "Consistency Analysis",
                      chart_con, 9, total_slides,
                      subtitle="TPS range (min-max) and standard deviation -- wider range = less predictable")

    # 10. Head-to-Head
    build_chart_slide(prs, "Head-to-Head: Qwen3-Coder vs All Models",
                      chart_h2h, 10, total_slides,
                      subtitle="Radar comparison across all 7 task categories")

    # 11. Takeaways
    build_takeaways_slide(prs, 11, total_slides)

    # Save
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_FILE))
    print(f"\nPresentation saved to: {OUTPUT_FILE}")
    print(f"File size: {OUTPUT_FILE.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
