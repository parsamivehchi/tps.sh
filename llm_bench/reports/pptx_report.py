"""Generate a polished, visually rich PowerPoint presentation."""

import json
from io import BytesIO
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from rich.console import Console

from llm_bench.config import (
    REPORTS_DIR, RESULTS_DIR, SCORED_DIR,
    MODELS, CATEGORY_LABELS, QUALITY_WEIGHTS,
)
from llm_bench.reports.charts import (
    speed_chart, latency_chart, quality_chart, cost_chart,
    category_heatmap, combined_speed_quality, model_comparison_radar,
)

console = Console()

# ── Design constants ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BRAND_DARK = RGBColor(0x0F, 0x17, 0x2A)
BRAND_BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xEC, 0x48, 0x99)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)


def _add_bg_rect(slide, color=BRAND_DARK):
    """Add full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send to back
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape


def _add_accent_bar(slide, top=Inches(1.8), color=BRAND_BLUE):
    """Add a thin accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(1.2), Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def _add_bullet_list(slide, left, top, width, height, items,
                     font_size=16, color=DARK_TEXT, spacing=Pt(8)):
    """Add a multi-line bullet list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def _add_chart_to_slide(slide, buf, left=Inches(0.4), top=Inches(1.6),
                         width=Inches(12.5)):
    """Add a chart image to a slide."""
    if buf is None or buf.getbuffer().nbytes == 0:
        return
    slide.shapes.add_picture(buf, left, top, width=width)


def _add_metric_card(slide, left, top, width, height, label, value, subtitle, accent_color):
    """Add a styled metric card shape."""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(1)

    # Accent top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Label
    _add_text_box(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.3),
                  label, font_size=10, color=MID_GRAY, bold=False)
    # Value
    _add_text_box(slide, left + Inches(0.15), top + Inches(0.4), width - Inches(0.3), Inches(0.5),
                  value, font_size=28, color=accent_color, bold=True)
    # Subtitle
    _add_text_box(slide, left + Inches(0.15), top + Inches(0.9), width - Inches(0.3), Inches(0.3),
                  subtitle, font_size=11, color=MID_GRAY, bold=False)


def generate_pptx_report(run_id: str | None = None) -> Path:
    """Generate polished LLM-Bench PowerPoint."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # Blank layout

    # ══════════════ Slide 1: Title ══════════════
    slide = prs.slides.add_slide(blank)
    _add_bg_rect(slide, BRAND_DARK)

    # Decorative accent bars
    for i, (c, w) in enumerate([(BRAND_BLUE, 2.0), (GREEN, 1.5), (AMBER, 1.0)]):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.5 + i * 0.15), Inches(w), Inches(0.04)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()

    _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(1.2),
                  "LLM-Bench", font_size=52, color=WHITE, bold=True)
    _add_text_box(slide, Inches(0.8), Inches(3.0), Inches(10), Inches(0.8),
                  "Local vs Cloud Model Benchmarking", font_size=24, color=RGBColor(0xCB, 0xD5, 0xE1))

    subtitle_text = f"Run: {run_id}" if run_id else "Methodology Overview"
    subtitle_text += f"  |  {datetime.now().strftime('%B %Y')}"
    _add_text_box(slide, Inches(0.8), Inches(4.0), Inches(10), Inches(0.5),
                  subtitle_text, font_size=14, color=MID_GRAY)

    _add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8),
                  "4 local open-source models (Ollama, M2 Max 32GB) vs 3 cloud Claude models (Anthropic API)\n"
                  "21 prompts across 7 categories  |  147 total test executions  |  LLM-as-judge quality scoring",
                  font_size=13, color=MID_GRAY)

    # ══════════════ Slide 2: Agenda ══════════════
    slide = prs.slides.add_slide(blank)
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                  "Agenda", font_size=32, color=BRAND_DARK, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    agenda_items = [
        "Models Under Test — 4 local + 3 cloud models",
        "Test Design — 7 categories, 21 prompts, what they measure",
        "Scoring Methodology — LLM-as-judge with bias disclosure",
        "Results — Speed, latency, quality, and cost charts",
        "Key Findings — Winners by category and use case",
        "Recommendations — Which model for which task",
    ]
    _add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
                     [f"  {item}" for item in agenda_items],
                     font_size=17, color=DARK_TEXT, spacing=Pt(14))

    # ══════════════ Slide 3: Models — split layout ══════════════
    slide = prs.slides.add_slide(blank)
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                  "Models Under Test", font_size=32, color=BRAND_DARK, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    # Left: Local models
    _add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4),
                  "LOCAL (Ollama on M2 Max, 32GB)", font_size=12, color=GREEN, bold=True)

    local_models = [m for m in MODELS if m.provider == "ollama"]
    for i, m in enumerate(local_models):
        y = Inches(2.2) + Inches(i * 1.1)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.8), y, Inches(5.5), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.fill.background()
        _add_text_box(slide, Inches(1.0), y + Inches(0.08), Inches(5), Inches(0.35),
                      m.name, font_size=15, color=DARK_TEXT, bold=True)
        _add_text_box(slide, Inches(1.0), y + Inches(0.42), Inches(5), Inches(0.35),
                      f"{m.model_type}  |  {m.model_id}  |  Free", font_size=11, color=MID_GRAY)

    # Right: Cloud models
    _add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.4),
                  "CLOUD (Anthropic API)", font_size=12, color=BRAND_BLUE, bold=True)

    cloud_models = [m for m in MODELS if m.provider == "anthropic"]
    for i, m in enumerate(cloud_models):
        y = Inches(2.2) + Inches(i * 1.1)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(7.0), y, Inches(5.5), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
        card.line.fill.background()
        _add_text_box(slide, Inches(7.2), y + Inches(0.08), Inches(5), Inches(0.35),
                      m.name, font_size=15, color=DARK_TEXT, bold=True)
        _add_text_box(slide, Inches(7.2), y + Inches(0.42), Inches(5), Inches(0.35),
                      f"{m.model_type}  |  ${m.cost_input:.0f}/${m.cost_output:.0f} per 1M tokens",
                      font_size=11, color=MID_GRAY)

    # ══════════════ Slide 4: Test Categories ══════════════
    slide = prs.slides.add_slide(blank)
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                  "7 Test Categories  ×  3 Prompts  =  21 Tests per Model",
                  font_size=28, color=BRAND_DARK, bold=True)
    _add_accent_bar(slide, Inches(1.15))

    cat_colors = [GREEN, BRAND_BLUE, AMBER, ROSE, PURPLE, RGBColor(0x06, 0xB6, 0xD4), RGBColor(0xF9, 0x73, 0x16)]
    cat_descs_short = [
        "Write code from specs",
        "Find bugs, trace logic",
        "Quick conversions",
        "Deep analysis & design",
        "Restructure code",
        "Teach & explain",
        "Tool use & agents",
    ]

    cols = 4
    card_w = Inches(2.8)
    card_h = Inches(1.5)
    start_x = Inches(0.6)
    start_y = Inches(1.7)
    gap_x = Inches(0.25)
    gap_y = Inches(0.2)

    for i, (cat_id, label) in enumerate(CATEGORY_LABELS.items()):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_GRAY
        card.line.width = Pt(1)

        # Color bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = cat_colors[i % len(cat_colors)]
        bar.line.fill.background()

        _add_text_box(slide, x + Inches(0.2), y + Inches(0.15), card_w - Inches(0.3), Inches(0.4),
                      label, font_size=14, color=DARK_TEXT, bold=True)
        _add_text_box(slide, x + Inches(0.2), y + Inches(0.55), card_w - Inches(0.3), Inches(0.5),
                      cat_descs_short[i] if i < len(cat_descs_short) else "",
                      font_size=11, color=MID_GRAY)
        _add_text_box(slide, x + Inches(0.2), y + Inches(1.0), card_w - Inches(0.3), Inches(0.3),
                      "3 prompts", font_size=10, color=cat_colors[i % len(cat_colors)], bold=True)

    # Total badge
    _add_text_box(slide, Inches(4), Inches(5.8), Inches(5), Inches(0.5),
                  "Total per full run: 7 models × 21 prompts = 147 executions",
                  font_size=14, color=BRAND_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # ══════════════ Slide 5: Scoring Methodology ══════════════
    slide = prs.slides.add_slide(blank)
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                  "Scoring Methodology", font_size=32, color=BRAND_DARK, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    # Three score dimension cards
    dims = [
        ("Correctness", "40%", "Technical accuracy.\nDoes the code work?\nAre facts correct?", BRAND_BLUE),
        ("Completeness", "35%", "Requirement coverage.\nAll parts addressed?\nEdge cases handled?", GREEN),
        ("Clarity", "25%", "Organization.\nWell-structured?\nEasy to follow?", AMBER),
    ]
    for i, (name, weight, desc, color) in enumerate(dims):
        x = Inches(0.8) + Inches(i * 4.0)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, Inches(1.8), Inches(3.6), Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_GRAY
        card.line.width = Pt(1)

        # Top bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(3.6), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        _add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(3.2), Inches(0.4),
                      name, font_size=18, color=DARK_TEXT, bold=True)
        _add_text_box(slide, x + Inches(0.2), Inches(2.4), Inches(3.2), Inches(0.4),
                      weight, font_size=28, color=color, bold=True)
        _add_text_box(slide, x + Inches(0.2), Inches(3.0), Inches(3.2), Inches(1.2),
                      desc, font_size=12, color=MID_GRAY)

    # Formula
    _add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.5),
                  "Weighted Score = 0.40 × Correctness + 0.35 × Completeness + 0.25 × Clarity",
                  font_size=16, color=BRAND_DARK, bold=True, align=PP_ALIGN.CENTER)

    # Bias note
    bias_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(1.5), Inches(5.9), Inches(10), Inches(0.8))
    bias_shape.fill.solid()
    bias_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xEB)
    bias_shape.line.color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    _add_text_box(slide, Inches(1.7), Inches(5.95), Inches(9.6), Inches(0.7),
                  "Bias Note: Claude Sonnet 4.6 judges all outputs including Claude family models. "
                  "Claude scores are flagged (*) in results. Interpret with care.",
                  font_size=12, color=RGBColor(0x92, 0x40, 0x0E))

    # ══════════════ Results Slides (if run data available) ══════════════
    if run_id:
        # Highlight metrics slide
        results_file = RESULTS_DIR / run_id / "results.json"
        analysis_file = SCORED_DIR / run_id / "analysis.json"

        if analysis_file.exists():
            analysis = json.loads(analysis_file.read_text())
            rankings = analysis.get("model_rankings", [])

            if rankings:
                slide = prs.slides.add_slide(blank)
                _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                              "Key Metrics at a Glance", font_size=32, color=BRAND_DARK, bold=True)
                _add_accent_bar(slide, Inches(1.2))

                fastest = min(rankings, key=lambda x: x.get("ttft_ms", float("inf")))
                best_tps = max(rankings, key=lambda x: x.get("tokens_per_sec", 0))

                cards = [
                    ("FASTEST TTFT", f"{fastest['ttft_ms']:.0f}ms", fastest['model_name'], PURPLE),
                    ("HIGHEST TPS", f"{best_tps['tokens_per_sec']:.1f}", best_tps['model_name'], GREEN),
                ]

                if analysis.get("has_quality_scores"):
                    best_q = max(rankings, key=lambda x: x.get("score_weighted", 0))
                    cards.append(("BEST QUALITY", f"{best_q['score_weighted']:.1f}/10", best_q['model_name'], BRAND_BLUE))

                paid = [r for r in rankings if r.get("cost_usd", 0) > 0]
                if paid:
                    cheapest = min(paid, key=lambda x: x["cost_usd"])
                    cards.append(("LOWEST COST", f"${cheapest['cost_usd']:.4f}", cheapest['model_name'], AMBER))

                card_w = Inches(2.8)
                for i, (label, value, sub, color) in enumerate(cards):
                    x = Inches(0.8) + Inches(i * 3.1)
                    _add_metric_card(slide, x, Inches(1.8), card_w, Inches(1.4), label, value, sub, color)

        # Chart slides
        chart_slides = [
            ("Generation Speed", speed_chart),
            ("Time to First Token", latency_chart),
            ("Quality Sub-Scores", quality_chart),
            ("Quality Radar", model_comparison_radar),
            ("Category Heatmap", category_heatmap),
            ("Speed vs Quality", combined_speed_quality),
            ("API Costs", cost_chart),
        ]

        for title, chart_fn in chart_slides:
            buf = chart_fn(run_id)
            if buf is not None and buf.getbuffer().nbytes > 0:
                slide = prs.slides.add_slide(blank)
                _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
                              title, font_size=28, color=BRAND_DARK, bold=True)
                _add_chart_to_slide(slide, buf, left=Inches(0.4), top=Inches(1.2), width=Inches(12.5))

    else:
        slide = prs.slides.add_slide(blank)
        _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                      "Results", font_size=32, color=BRAND_DARK, bold=True)
        _add_accent_bar(slide, Inches(1.2))
        _add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4), [
            "Charts and data will appear after running benchmarks",
            "",
            "Step 1: python -m llm_bench run",
            "Step 2: python -m llm_bench judge <run_id>",
            "Step 3: python -m llm_bench analyze <run_id>",
            "Step 4: python -m llm_bench report <run_id>",
        ], font_size=16, color=DARK_TEXT)

    # ══════════════ Recommendations Slide ══════════════
    slide = prs.slides.add_slide(blank)
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                  "Which Model for Which Task?", font_size=32, color=BRAND_DARK, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    recs = [
        ("Batch processing / code review", "Local models — zero cost at any scale", GREEN),
        ("Interactive assistant / chat", "Lowest TTFT model for perceived speed", PURPLE),
        ("Production-critical code", "Highest quality model (usually cloud)", BRAND_BLUE),
        ("Quick prototyping", "Fast local model with adequate quality (>6/10)", AMBER),
        ("Cost-sensitive deployment", "Best quality-per-dollar cloud model", ROSE),
    ]

    for i, (scenario, rec, color) in enumerate(recs):
        y = Inches(1.7) + Inches(i * 1.0)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.8), y, Inches(0.06), Inches(0.7))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        _add_text_box(slide, Inches(1.1), y, Inches(5), Inches(0.4),
                      scenario, font_size=15, color=DARK_TEXT, bold=True)
        _add_text_box(slide, Inches(1.1), y + Inches(0.35), Inches(11), Inches(0.35),
                      rec, font_size=13, color=MID_GRAY)

    # ══════════════ Final Slide ══════════════
    slide = prs.slides.add_slide(blank)
    _add_bg_rect(slide, BRAND_DARK)

    for i, (c, w) in enumerate([(BRAND_BLUE, 2.0), (GREEN, 1.5), (AMBER, 1.0)]):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.6 + i * 0.15), Inches(w), Inches(0.04)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()

    _add_text_box(slide, Inches(0.8), Inches(2.2), Inches(10), Inches(1),
                  "Thank You", font_size=44, color=WHITE, bold=True)
    _add_text_box(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.5),
                  "LLM-Bench — Open-source model benchmarking toolkit",
                  font_size=16, color=MID_GRAY)
    _add_text_box(slide, Inches(0.8), Inches(5.0), Inches(10), Inches(0.8),
                  "python -m llm_bench run  |  judge  |  analyze  |  report  |  export",
                  font_size=13, color=RGBColor(0x94, 0xA3, 0xB8))

    # ── Save ──
    output_path = REPORTS_DIR / f"llm_bench_presentation{'_' + run_id if run_id else ''}.pptx"
    prs.save(str(output_path))
    console.print(f"[green]PowerPoint saved: {output_path}[/]")
    return output_path
