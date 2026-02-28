"""Generate the On-Premise LLM Infrastructure PowerPoint presentation."""

from datetime import datetime
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from rich.console import Console

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from llm_bench.config import REPORTS_DIR
from llm_bench.analysis.hardware_analysis import (
    HARDWARE_SPECS, MODEL_SPECS, CLUSTER_CONFIGS,
    ACTUAL_BENCHMARKS, CLOUD_PRICING, FRAMEWORK_COMPARISON,
    QUANTIZATION_INFO,
    projected_tps, projected_tps_moe, cost_per_million_tokens,
    project_all_hardware,
)

console = Console()

# ── Design constants ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BRAND_DARK = RGBColor(0x0F, 0x17, 0x2A)
BRAND_BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xEC, 0x48, 0x99)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)


def _add_bg(slide, color=BRAND_DARK):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def _add_accent_bar(slide, top=Inches(1.8), color=BRAND_BLUE, left=Inches(0.8), width=Inches(1.2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _text(slide, left, top, width, height, text, size=18,
          color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return txBox


def _bullets(slide, left, top, width, height, items, size=16, color=DARK_TEXT, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
    return txBox


def _metric_card(slide, left, top, width, height, label, value, subtitle, accent):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(1)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    _text(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.3),
          label, size=10, color=MID_GRAY)
    _text(slide, left + Inches(0.15), top + Inches(0.4), width - Inches(0.3), Inches(0.5),
          value, size=28, color=accent, bold=True)
    _text(slide, left + Inches(0.15), top + Inches(0.9), width - Inches(0.3), Inches(0.3),
          subtitle, size=11, color=MID_GRAY)


def _add_chart(slide, buf, left=Inches(0.4), top=Inches(1.6), width=Inches(12.5)):
    if buf is None or buf.getbuffer().nbytes == 0:
        return
    slide.shapes.add_picture(buf, left, top, width=width)


# ── Chart generators ──

def _bandwidth_chart() -> BytesIO:
    fig, ax = plt.subplots(figsize=(12, 4))

    names = [hw.chip for hw in HARDWARE_SPECS]
    bws = [hw.memory_bandwidth_gbs for hw in HARDWARE_SPECS]
    colors = ["#10b981", "#3b82f6", "#8b5cf6", "#8b5cf6", "#ec4899", "#f59e0b"]

    bars = ax.barh(names, bws, color=colors, height=0.5)
    for bar, bw in zip(bars, bws):
        ax.text(bar.get_width() + 8, bar.get_y() + bar.get_height() / 2,
                f"{bw} GB/s", va="center", fontsize=11, fontweight="bold")

    ax.set_xlabel("Memory Bandwidth (GB/s)", fontsize=11)
    ax.invert_yaxis()
    ax.set_xlim(0, max(bws) * 1.18)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout()

    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def _tps_projection_chart() -> BytesIO:
    model = next(m for m in MODEL_SPECS if "qwen3-coder" in m.name)
    projections = project_all_hardware(model, "q4")

    fig, ax = plt.subplots(figsize=(12, 5.5))

    names = [p["hardware"].replace(" (current)", "\n(yours)").replace(" (upcoming)", "\n(upcoming)")
             for p in projections]
    tps = [p["projected_tps"] for p in projections]
    colors = []
    for p in projections:
        if p["form_factor"] == "cluster":
            colors.append("#f59e0b")
        elif "(current)" in p["hardware"]:
            colors.append("#10b981")
        else:
            colors.append("#3b82f6")

    bars = ax.barh(names, tps, color=colors, height=0.55)
    ax.axvline(x=48.8, color="#ef4444", linestyle="--", linewidth=2, alpha=0.6)
    ax.text(49.5, len(names) - 0.5, "Your actual: 48.8 TPS", fontsize=9, color="#ef4444")

    for bar, t in zip(bars, tps):
        ax.text(bar.get_width() + 1, bar.get_y() + bar.get_height() / 2,
                f"{t:.0f}", va="center", fontsize=10, fontweight="bold")

    ax.set_xlabel("Projected TPS (qwen3-coder Q4)", fontsize=11)
    ax.invert_yaxis()
    ax.set_xlim(0, max(tps) * 1.15)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout()

    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def _cost_chart() -> BytesIO:
    fig, ax = plt.subplots(figsize=(12, 4))

    items = []
    for hw in HARDWARE_SPECS:
        if hw.price_usd == 0:
            continue
        tps = projected_tps(hw.memory_bandwidth_gbs, 15, 0.6)
        cost = cost_per_million_tokens(hw.price_usd, tps)
        items.append((hw.chip, cost, "#10b981"))

    items.append(("Claude Haiku", 5.0, "#06b6d4"))
    items.append(("Claude Sonnet", 15.0, "#8b5cf6"))
    items.append(("Claude Opus", 25.0, "#ec4899"))

    names = [i[0] for i in items]
    costs = [i[1] for i in items]
    colors = [i[2] for i in items]

    bars = ax.bar(names, costs, color=colors, width=0.55)
    for bar, c in zip(bars, costs):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3,
                f"${c:.2f}", ha="center", fontsize=10, fontweight="bold")

    ax.set_ylabel("$/1M Output Tokens", fontsize=11)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.xticks(rotation=25, ha="right", fontsize=10)
    plt.tight_layout()

    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def _cluster_scaling_chart() -> BytesIO:
    fig, ax = plt.subplots(figsize=(12, 4))

    nodes = [1, 2, 3, 4, 5, 6]
    bw_per_node = 273  # M4 Pro
    theoretical = [n * bw_per_node / 15 * 0.6 for n in nodes]
    actual_eff = [1.0, 0.85, 0.80, 0.75, 0.72, 0.70]
    realistic = [t * e for t, e in zip(theoretical, actual_eff)]

    ax.plot(nodes, theoretical, "--", color="#3b82f6", linewidth=2, label="Theoretical (linear)")
    ax.plot(nodes, realistic, "-o", color="#10b981", linewidth=2, markersize=8, label="Realistic (with overhead)")
    ax.fill_between(nodes, realistic, alpha=0.1, color="#10b981")

    for n, r in zip(nodes, realistic):
        ax.annotate(f"{r:.0f} TPS", (n, r), textcoords="offset points",
                    xytext=(10, 10), fontsize=9, fontweight="bold", color="#10b981")

    ax.set_xlabel("Number of M4 Pro Mac Minis", fontsize=11)
    ax.set_ylabel("Projected TPS (qwen3-coder Q4)", fontsize=11)
    ax.legend(fontsize=10)
    ax.set_xticks(nodes)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout()

    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


# ── Main generator ──

def generate_hardware_pptx() -> Path:
    """Generate the On-Premise LLM Infrastructure PowerPoint."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ══════════════ Slide 1: Title ══════════════
    slide = prs.slides.add_slide(blank)
    _add_bg(slide)

    for i, (c, w) in enumerate([(BRAND_BLUE, 2.5), (GREEN, 2.0), (AMBER, 1.5)]):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.9 + i * 0.15), Inches(w), Inches(0.04))
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()

    _text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(1.2),
          "On-Premise LLM Infrastructure", size=48, color=WHITE, bold=True)
    _text(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.8),
          "A Hardware Guide for Apple Silicon", size=24, color=RGBColor(0xCB, 0xD5, 0xE1))
    _text(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
          f"LLM-Bench Phase 2  |  {datetime.now().strftime('%B %Y')}", size=14, color=MID_GRAY)
    _text(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8),
          "Based on 147 real benchmarks across 7 models on M2 Max 32GB\n"
          "Memory bandwidth is the bottleneck — not GPU cores, not RAM size",
          size=13, color=MID_GRAY)

    # ══════════════ Slide 2: The Question ══════════════
    slide = prs.slides.add_slide(blank)
    _text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
          "What's Limiting Your Local LLMs?", size=32, color=DARK_TEXT, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    questions = [
        "Why is qwen3-coder only doing 48.8 tokens/sec?",
        "Would a 64GB MacBook make it faster?",
        "Would a better GPU help?",
        "How do I scale beyond one machine?",
    ]
    _bullets(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(3),
             [f"  {q}" for q in questions], size=20, color=DARK_TEXT, spacing=Pt(18))

    # Answer box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
    box.line.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)

    _text(slide, Inches(1.0), Inches(5.0), Inches(11), Inches(0.4),
          "The Answer: Memory Bandwidth", size=22, color=BRAND_BLUE, bold=True)
    _text(slide, Inches(1.0), Inches(5.5), Inches(11), Inches(0.8),
          "LLM inference is memory-bound, not compute-bound. Every token requires loading\n"
          "ALL model weights from memory. The speed of that transfer is your bottleneck.",
          size=14, color=RGBColor(0x1E, 0x40, 0x8A))

    # ══════════════ Slide 3: The Answer — Bandwidth Diagram ══════════════
    slide = prs.slides.add_slide(blank)
    _text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
          "Memory Bandwidth = The Bottleneck", size=32, color=DARK_TEXT, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    # Formula box
    formula_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(1.5), Inches(1.8), Inches(10), Inches(1.2))
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = BRAND_DARK
    formula_box.line.fill.background()

    _text(slide, Inches(1.8), Inches(1.9), Inches(9.5), Inches(0.5),
          "TPS = Memory Bandwidth (GB/s) / Model Size (GB) x Efficiency", size=22, color=WHITE, bold=True,
          align=PP_ALIGN.CENTER)
    _text(slide, Inches(1.8), Inches(2.45), Inches(9.5), Inches(0.4),
          "Real-world efficiency: 50-70% of theoretical maximum", size=13, color=MID_GRAY,
          align=PP_ALIGN.CENTER)

    # Pipe analogy — three cards
    analogies = [
        ("Bandwidth", "= Pipe Diameter", "How fast data moves\nfrom memory to GPU", BRAND_BLUE),
        ("Model Size", "= Water Volume", "How much data must flow\nfor each token", GREEN),
        ("Efficiency", "= Pipe Friction", "Framework overhead,\nKV cache, OS usage", AMBER),
    ]
    for i, (title, subtitle, desc, color) in enumerate(analogies):
        x = Inches(0.8) + Inches(i * 4.1)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, Inches(3.5), Inches(3.8), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_GRAY

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.5), Inches(3.8), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        _text(slide, x + Inches(0.2), Inches(3.7), Inches(3.4), Inches(0.4),
              title, size=18, color=DARK_TEXT, bold=True)
        _text(slide, x + Inches(0.2), Inches(4.1), Inches(3.4), Inches(0.4),
              subtitle, size=14, color=color, bold=True)
        _text(slide, x + Inches(0.2), Inches(4.6), Inches(3.4), Inches(0.8),
              desc, size=12, color=MID_GRAY)

    # Validation
    _text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
          "Validated: M2 Max predicted 55 TPS for qwen3-coder, actual = 48.8 TPS (11% delta)",
          size=13, color=BRAND_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # ══════════════ Slide 4: Your Results ══════════════
    slide = prs.slides.add_slide(blank)
    _text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
          "Your Phase 1 Results (M2 Max 32GB)", size=32, color=DARK_TEXT, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    # Metric cards for key results
    cards_data = [
        ("FASTEST LOCAL", "48.8 TPS", "qwen3-coder (30B MoE)", GREEN),
        ("BEST QUALITY", "8.65/10", "Claude Opus 4.6", BRAND_BLUE),
        ("FASTEST OVERALL", "169.7 TPS", "Claude Haiku 4.5", PURPLE),
        ("TOTAL COST", "$3.95", "147 benchmarks", AMBER),
    ]
    for i, (label, value, sub, color) in enumerate(cards_data):
        x = Inches(0.6) + Inches(i * 3.15)
        _metric_card(slide, x, Inches(1.6), Inches(2.9), Inches(1.4), label, value, sub, color)

    # Results table
    _text(slide, Inches(0.8), Inches(3.3), Inches(11), Inches(0.4),
          "All 7 Models Compared", size=16, color=DARK_TEXT, bold=True)

    results = [
        ("Claude Haiku 4.5", "Cloud", "169.7", "500ms", "8.25", "$0.28"),
        ("Claude Sonnet 4.6", "Cloud", "77.7", "1.0s", "8.59", "$0.93"),
        ("Claude Opus 4.6", "Cloud", "76.6", "1.8s", "8.65", "$1.48"),
        ("qwen3-coder", "Local", "48.8", "1.1s", "7.48", "Free"),
        ("qwen2.5-coder:14b", "Local", "15.6", "1.5s", "6.64", "Free"),
        ("deepseek-r1:14b", "Local", "14.6", "70.2s", "5.89", "Free"),
        ("glm-4.7-flash", "Local", "10.2", "54.8s", "5.30", "Free"),
    ]

    # Build table as text boxes (more control than pptx tables)
    headers = ["Model", "Type", "TPS", "TTFT", "Quality", "Cost"]
    col_x = [Inches(0.8), Inches(4.2), Inches(5.8), Inches(7.2), Inches(8.8), Inches(10.5)]
    col_w = [Inches(3.4), Inches(1.6), Inches(1.4), Inches(1.6), Inches(1.7), Inches(1.5)]

    header_y = Inches(3.7)
    for j, h in enumerate(headers):
        _text(slide, col_x[j], header_y, col_w[j], Inches(0.3),
              h, size=10, color=WHITE, bold=True)

    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.7), header_y, Inches(11.5), Inches(0.32))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = BRAND_DARK
    header_bg.line.fill.background()
    sp = header_bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

    for i, row in enumerate(results):
        y = Inches(4.1) + Inches(i * 0.4)
        for j, val in enumerate(row):
            _text(slide, col_x[j], y, col_w[j], Inches(0.35),
                  val, size=11, color=DARK_TEXT)

    # ══════════════ Slide 5: Hardware Lineup ══════════════
    slide = prs.slides.add_slide(blank)
    _text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
          "Apple Silicon Hardware Lineup", size=32, color=DARK_TEXT, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    hw_colors = [GREEN, BRAND_BLUE, PURPLE, PURPLE, ROSE, AMBER]
    for i, hw in enumerate(HARDWARE_SPECS):
        col = i % 3
        row = i // 3
        x = Inches(0.6) + Inches(col * 4.15)
        y = Inches(1.6) + Inches(row * 2.9)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, y, Inches(3.9), Inches(2.6))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_GRAY

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.9), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hw_colors[i]
        bar.line.fill.background()

        short_name = hw.name.split("(")[0].strip()
        _text(slide, x + Inches(0.15), y + Inches(0.15), Inches(3.6), Inches(0.35),
              short_name, size=14, color=DARK_TEXT, bold=True)

        price_str = "Owned" if hw.price_usd == 0 else f"${hw.price_usd:,}"
        _text(slide, x + Inches(0.15), y + Inches(0.5), Inches(3.6), Inches(0.3),
              f"{hw.chip}  |  {price_str}", size=11, color=hw_colors[i], bold=True)

        specs = (f"{hw.memory_bandwidth_gbs} GB/s bandwidth\n"
                 f"{hw.max_ram_gb} GB max RAM\n"
                 f"{hw.gpu_cores} GPU cores")
        _text(slide, x + Inches(0.15), y + Inches(0.9), Inches(3.6), Inches(1.0),
              specs, size=11, color=MID_GRAY)

        note = hw.notes if len(hw.notes) < 55 else hw.notes[:52] + "..."
        _text(slide, x + Inches(0.15), y + Inches(2.0), Inches(3.6), Inches(0.4),
              note, size=9, color=RGBColor(0x94, 0xA3, 0xB8))

    # ══════════════ Slide 6: TPS Projections Chart ══════════════
    slide = prs.slides.add_slide(blank)
    _text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
          "Projected TPS: qwen3-coder Across Hardware", size=28, color=DARK_TEXT, bold=True)
    _add_chart(slide, _tps_projection_chart(), top=Inches(1.2))

    # ══════════════ Slide 7: RAM vs Speed ══════════════
    slide = prs.slides.add_slide(blank)
    _text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
          "More RAM \u2260 Faster Inference", size=32, color=DARK_TEXT, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    # Two columns
    # Left: RAM determines what fits
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
    left_box.line.color.rgb = RGBColor(0xBB, 0xF7, 0xD0)

    _text(slide, Inches(0.9), Inches(2.0), Inches(5.2), Inches(0.4),
          "RAM = What You Can Run", size=20, color=GREEN, bold=True)
    _bullets(slide, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.5), [
        "32GB: Models up to ~25GB (14B Q4, 30B MoE Q4)",
        "64GB: Models up to ~55GB (70B Q4)",
        "128GB: Models up to ~115GB (70B Q8, 123B Q4)",
        "512GB: Models up to ~450GB (405B Q4, 671B MoE Q4)",
    ], size=13, color=DARK_TEXT, spacing=Pt(14))

    # Right: Bandwidth determines speed
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
    right_box.line.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)

    _text(slide, Inches(7.1), Inches(2.0), Inches(5.2), Inches(0.4),
          "Bandwidth = How Fast It Runs", size=20, color=BRAND_BLUE, bold=True)
    _bullets(slide, Inches(7.1), Inches(2.5), Inches(5.2), Inches(3.5), [
        "273 GB/s (M4 Pro): ~33 TPS for qwen3-coder",
        "400 GB/s (M2 Max): ~49 TPS (your current)",
        "546 GB/s (M4 Max): ~67 TPS (+37%!)",
        "819 GB/s (M3 Ultra): ~100 TPS (+105%!)",
    ], size=13, color=DARK_TEXT, spacing=Pt(14))

    # Warning callout
    warn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1.5), Inches(6.5), Inches(10), Inches(0.7))
    warn.fill.solid()
    warn.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xEB)
    warn.line.color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    _text(slide, Inches(1.7), Inches(6.55), Inches(9.6), Inches(0.6),
          "A 64GB M4 Pro Mac Mini (273 GB/s) will be SLOWER than your M2 Max (400 GB/s) for the same model!",
          size=14, color=RGBColor(0x92, 0x40, 0x0E), bold=True, align=PP_ALIGN.CENTER)

    # ══════════════ Slide 8: MoE Architecture ══════════════
    slide = prs.slides.add_slide(blank)
    _text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
          "MoE: Why qwen3-coder Punches Above Its Weight", size=28, color=DARK_TEXT, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    # Dense vs MoE comparison cards
    dense_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.6), Inches(1.6), Inches(5.8), Inches(3))
    dense_card.fill.solid()
    dense_card.fill.fore_color.rgb = WHITE
    dense_card.line.color.rgb = LIGHT_GRAY

    _text(slide, Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.4),
          "Dense Model (14B)", size=18, color=ROSE, bold=True)
    _bullets(slide, Inches(0.9), Inches(2.3), Inches(5.2), Inches(2), [
        "All 14B parameters activate every token",
        "~7.4 GB in Q4 quantization",
        "TPS = 400 / 7.4 x 0.6 = 32 (predicted)",
        "Actual: 15.6 TPS (reasoning overhead)",
    ], size=13, color=DARK_TEXT, spacing=Pt(10))

    moe_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(6.8), Inches(1.6), Inches(5.8), Inches(3))
    moe_card.fill.solid()
    moe_card.fill.fore_color.rgb = WHITE
    moe_card.line.color.rgb = LIGHT_GRAY

    _text(slide, Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4),
          "MoE Model (30B total, 8B active)", size=18, color=GREEN, bold=True)
    _bullets(slide, Inches(7.1), Inches(2.3), Inches(5.2), Inches(2), [
        "Only 8B of 30B params activate per token",
        "~15.3 GB total, but ~4 GB active",
        "TPS = 400 / 4 x 0.55 = 55 (predicted)",
        "Actual: 48.8 TPS (close to prediction!)",
    ], size=13, color=DARK_TEXT, spacing=Pt(10))

    # Key insight
    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(1.5), Inches(5.0), Inches(10), Inches(1.5))
    insight.fill.solid()
    insight.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
    insight.line.color.rgb = RGBColor(0xBB, 0xF7, 0xD0)
    _text(slide, Inches(1.8), Inches(5.1), Inches(9.4), Inches(0.4),
          "MoE = Best of Both Worlds", size=18, color=GREEN, bold=True)
    _text(slide, Inches(1.8), Inches(5.5), Inches(9.4), Inches(0.8),
          "30B total parameters = broad knowledge from diverse experts\n"
          "8B active parameters = speed equivalent to a much smaller model\n"
          "Tradeoff: Full 15GB must still fit in RAM even though only 4GB activates",
          size=13, color=DARK_TEXT)

    # ══════════════ Slide 9: Scaling Up ══════════════
    slide = prs.slides.add_slide(blank)
    _text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
          "Scaling Up: Mac Studio vs Mac Mini Cluster", size=28, color=DARK_TEXT, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    # Two paths
    _text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
          "Path A: Single Powerful Node", size=16, color=PURPLE, bold=True)
    _bullets(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5), [
        "M4 Max Mac Studio ($4,499)",
        "546 GB/s bandwidth, 128GB RAM",
        "~67 TPS for qwen3-coder",
        "No network overhead, simple setup",
        "Best for: single-user, latency-sensitive",
    ], size=13, color=DARK_TEXT, spacing=Pt(8))

    _text(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.4),
          "Path B: Multi-Node Cluster (exo)", size=16, color=AMBER, bold=True)
    _bullets(slide, Inches(6.8), Inches(2.1), Inches(5.5), Inches(2.5), [
        "2x M4 Pro Mac Mini ($3,598)",
        "546 GB/s aggregate, 96GB total RAM",
        "~56 TPS for qwen3-coder",
        "Scalable — add more nodes later",
        "Best for: teams, budget-conscious, big models",
    ], size=13, color=DARK_TEXT, spacing=Pt(8))

    # Cluster scaling chart
    _add_chart(slide, _cluster_scaling_chart(), top=Inches(4.5), left=Inches(0.5), width=Inches(12.3))

    # ══════════════ Slide 10: exo Distributed ══════════════
    slide = prs.slides.add_slide(blank)
    _text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
          "exo: Distributed Inference Across Macs", size=28, color=DARK_TEXT, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    _bullets(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4), [
        "How it works:",
        "  Model layers split across nodes",
        "  Activations passed via Thunderbolt 5",
        "  Automatic peer discovery on LAN",
        "  Compatible with Ollama & MLX models",
        "",
        "Requirements:",
        "  Thunderbolt 5 (120 Gbps) for best results",
        "  Same architecture models on all nodes",
        "  pip install exo on each Mac",
    ], size=14, color=DARK_TEXT, spacing=Pt(6))

    # Cluster configs
    _text(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.4),
          "Example Configurations", size=16, color=DARK_TEXT, bold=True)

    for i, c in enumerate(CLUSTER_CONFIGS):
        y = Inches(2.2) + Inches(i * 1.5)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(6.8), y, Inches(5.5), Inches(1.3))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.fill.background()

        _text(slide, Inches(7.0), y + Inches(0.1), Inches(5.1), Inches(0.3),
              c["name"], size=13, color=DARK_TEXT, bold=True)
        _text(slide, Inches(7.0), y + Inches(0.4), Inches(5.1), Inches(0.7),
              f"{c['aggregate_bandwidth_gbs']} GB/s  |  {c['aggregate_ram_gb']}GB RAM  |  ${c['price_usd']:,}\n"
              f"{c['interconnect']}  |  {c['scaling_efficiency']:.0%} efficiency",
              size=11, color=MID_GRAY)

    # ══════════════ Slide 11: Cost Analysis ══════════════
    slide = prs.slides.add_slide(blank)
    _text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
          "Cost: Local (Amortized) vs Cloud API", size=28, color=DARK_TEXT, bold=True)
    _add_chart(slide, _cost_chart(), top=Inches(1.0))

    _text(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.5),
          "Local cost assumes: 3-year lifespan, 25% utilization, $0.15/kWh, 30W power draw",
          size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)
    _text(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
          "At >5M tokens/month, local hardware pays for itself within 6-12 months vs Claude API",
          size=14, color=BRAND_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # ══════════════ Slide 12: Decision Matrix ══════════════
    slide = prs.slides.add_slide(blank)
    _text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
          "Decision Matrix: What Should You Buy?", size=28, color=DARK_TEXT, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    recs = [
        ("< $2K / Personal", "M4 Pro Mac Mini 48GB", "Budget friendly always-on LLM server", GREEN),
        ("$2-4K / Personal", "M4 Max MacBook Pro 64GB", "37% faster than your M2 Max, portable", BRAND_BLUE),
        ("$2-4K / Team", "2x M4 Pro Mac Mini (exo)", "Shared cluster, scalable, 96GB total", AMBER),
        ("$4-6K / Team", "M4 Max Mac Studio 128GB", "Desktop powerhouse, superior thermals", PURPLE),
        ("$6K+ / Lab", "M3/M4 Ultra Mac Studio", "512GB RAM for frontier models (405B+)", ROSE),
        ("$8K+ / Lab", "2x M4 Max Mac Studio (exo)", "Maximum performance, 256GB total", CYAN),
    ]

    for i, (scenario, hardware, reason, color) in enumerate(recs):
        y = Inches(1.6) + Inches(i * 0.9)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.8), y, Inches(0.06), Inches(0.65))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        _text(slide, Inches(1.1), y, Inches(3.2), Inches(0.35),
              scenario, size=14, color=DARK_TEXT, bold=True)
        _text(slide, Inches(4.5), y, Inches(4), Inches(0.35),
              hardware, size=14, color=color, bold=True)
        _text(slide, Inches(1.1), y + Inches(0.3), Inches(11), Inches(0.35),
              reason, size=12, color=MID_GRAY)

    # ══════════════ Slide 13: One-Script Setup ══════════════
    slide = prs.slides.add_slide(blank)
    _text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
          "One-Script Mac Setup", size=32, color=DARK_TEXT, bold=True)
    _add_accent_bar(slide, Inches(1.2))

    _text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5),
          "git clone <repo>/mac-setup && cd mac-setup && ./setup.sh",
          size=16, color=BRAND_BLUE, bold=True)

    setup_items = [
        "Homebrew + essential CLI tools (bat, eza, fzf, ripgrep, fd, jq, htop)",
        "Development tools (git, gh, python, node, docker)",
        "LLM tools (Ollama + 5 recommended models pulled automatically)",
        "AI coding (Claude Code, Cursor, VS Code)",
        "Terminal (Ghostty + Starship prompt + zoxide)",
        "macOS defaults (dev-friendly Dock, Finder, keyboard settings)",
        "Git + SSH key setup",
        "Shell configuration (zsh aliases, PATH, environment variables)",
    ]
    _bullets(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(4),
             [f"  {item}" for item in setup_items],
             size=14, color=DARK_TEXT, spacing=Pt(10))

    _text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
          "Idempotent — safe to re-run on an already-configured Mac",
          size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # ══════════════ Slide 14: Next Steps ══════════════
    slide = prs.slides.add_slide(blank)
    _add_bg(slide)

    for i, (c, w) in enumerate([(BRAND_BLUE, 2.5), (GREEN, 2.0), (AMBER, 1.5)]):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.6 + i * 0.15), Inches(w), Inches(0.04))
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()

    _text(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(1),
          "Next Steps", size=44, color=WHITE, bold=True)

    steps = [
        "1. Order hardware (when ready)",
        "2. Run setup.sh on new machine",
        "3. Run LLM-Bench with --hardware-metrics",
        "4. Compare actual vs projected TPS",
        "5. Share results with team",
    ]
    _bullets(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(2.5),
             steps, size=18, color=RGBColor(0xCB, 0xD5, 0xE1), spacing=Pt(12))

    _text(slide, Inches(0.8), Inches(6.5), Inches(10), Inches(0.5),
          "LLM-Bench Phase 2 — On-Premise Infrastructure Guide",
          size=13, color=MID_GRAY)

    # ── Save ──
    output_path = REPORTS_DIR / "on_premise_llm_infrastructure.pptx"
    prs.save(str(output_path))
    console.print(f"[green]Hardware PPTX saved: {output_path}[/]")
    return output_path
